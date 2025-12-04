import streamlit as st
import os
import json
import shutil
from pathlib import Path
from datetime import datetime, date, timedelta

from dotenv import load_dotenv
from openai import OpenAI
from pypdf import PdfReader
from pptx import Presentation
import io
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
import textwrap
import re

# -----------------------------------------
# Paths and basic storage helpers
# -----------------------------------------

BASE_DIR = Path(__file__).parent
DATA_DIR = BASE_DIR / "data"
COURSES_FILE = DATA_DIR / "courses.json"
COURSES_DIR = DATA_DIR / "courses"


def ensure_data_dirs():
    """Ensure the data directories and files exist."""
    DATA_DIR.mkdir(exist_ok=True)
    COURSES_DIR.mkdir(exist_ok=True)
    if not COURSES_FILE.exists():
        COURSES_FILE.write_text(json.dumps({"courses": []}, indent=2))


# Load environment variables and set up OpenAI
load_dotenv()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None


def load_courses():
    """Load list of courses from courses.json."""
    if not COURSES_FILE.exists():
        return []
    try:
        data = json.loads(COURSES_FILE.read_text())
        return data.get("courses", [])
    except json.JSONDecodeError:
        return []


def save_courses(courses):
    """Save list of courses to courses.json."""
    COURSES_FILE.write_text(json.dumps({"courses": courses}, indent=2))


def create_course(course_name: str):
    """Create a new course with a unique id."""
    courses = load_courses()
    # Simple id scheme: course_1, course_2, ...
    existing_ids = [c["id"] for c in courses]
    new_index = len(existing_ids) + 1
    new_id = f"course_{new_index}"

    new_course = {"id": new_id, "name": course_name}
    courses.append(new_course)
    save_courses(courses)

    # Create course directory + meta
    course_dir = COURSES_DIR / new_id
    course_dir.mkdir(parents=True, exist_ok=True)
    meta_path = course_dir / "meta.json"
    if not meta_path.exists():
        meta = {"id": new_id, "name": course_name, "files": []}
        meta_path.write_text(json.dumps(meta, indent=2))

    return new_course


def get_course_by_name(course_name: str):
    """Find a course dict by name."""
    for c in load_courses():
        if c["name"] == course_name:
            return c
    return None


def get_course_by_id(course_id: str):
    """Find a course dict by id."""
    for c in load_courses():
        if c["id"] == course_id:
            return c
    return None


def delete_course(course_id: str):
    """Delete a course from courses.json and remove its directory from disk."""
    courses = load_courses()
    courses = [c for c in courses if c["id"] != course_id]
    save_courses(courses)

    course_dir = get_course_dir(course_id)
    if course_dir.exists():
        shutil.rmtree(course_dir)


def get_course_dir(course_id: str) -> Path:
    return COURSES_DIR / course_id


def get_course_meta_path(course_id: str) -> Path:
    return get_course_dir(course_id) / "meta.json"


def load_course_meta(course_id: str):
    """Load meta.json for a course (name + files)."""
    meta_path = get_course_meta_path(course_id)
    if not meta_path.exists():
        course = get_course_by_id(course_id) or {"id": course_id, "name": "Unknown Course"}
        meta = {"id": course_id, "name": course.get("name", "Unknown Course"), "files": []}
        meta_path.write_text(json.dumps(meta, indent=2))
        return meta

    try:
        return json.loads(meta_path.read_text())
    except json.JSONDecodeError:
        return {"id": course_id, "name": "Unknown Course", "files": []}


def save_course_meta(course_id: str, meta: dict):
    """Save meta.json for a course."""
    meta_path = get_course_meta_path(course_id)
    meta_path.write_text(json.dumps(meta, indent=2))


def make_unique_filename(folder: Path, original_name: str) -> str:
    """
    If a file with original_name exists in folder, append a counter.
    e.g., Lecture1.pdf -> Lecture1_1.pdf, Lecture1_2.pdf, etc.
    """
    base = Path(original_name).stem
    ext = Path(original_name).suffix
    candidate = original_name
    counter = 1
    while (folder / candidate).exists():
        candidate = f"{base}_{counter}{ext}"
        counter += 1
    return candidate


def extract_text_from_file(file_path: Path) -> str:
    """
    Extract plain text from a course file.
    Currently supports: PDF (.pdf), PowerPoint (.pptx/.ppt), and text files (.txt/.md).
    """
    suffix = file_path.suffix.lower()
    try:
        if suffix == ".pdf":
            reader = PdfReader(str(file_path))
            texts = []
            for page in reader.pages:
                content = page.extract_text() or ""
                texts.append(content)
            return "\n\n".join(texts)

        elif suffix in [".pptx", ".ppt"]:
            prs = Presentation(str(file_path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n\n".join(texts)

        elif suffix in [".txt", ".md"]:
            return file_path.read_text(errors="ignore")

        else:
            # Unsupported type for now
            return ""
    except Exception:
        # If anything goes wrong parsing a file, just skip its content
        return ""


def generate_flashcards_from_text(text: str, num_cards: int = 20) -> list[dict]:
    """
    Call the OpenAI API to generate term/definition flashcards from the given text.
    Returns a list of dicts: [{"front": "...", "back": "..."}, ...]
    """
    if not client:
        raise RuntimeError("OpenAI client is not configured. Missing OPENAI_API_KEY.")

    # Keep the prompt reasonably sized
    MAX_CHARS = 12000
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]

    prompt = f"""
You are an assistant that creates concise, high-quality term/definition flashcards to help a student study.

Given the course material below, generate up to {num_cards} of the most important flashcards.
Each flashcard must follow this schema:

[
  {{
    "front": "Short term or question",
    "back": "Clear, student-friendly definition or explanation"
  }},
  ...
]

Rules:
- Focus on key concepts, definitions, formulas, and important distinctions.
- Avoid trivial details.
- The 'front' should be short, like a term or a brief question.
- The 'back' should be 1–4 sentences, clear and precise.
- Respond with ONLY valid JSON (a list of objects), no extra text.

COURSE MATERIAL START
{text}
COURSE MATERIAL END
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You generate helpful, accurate term/definition flashcards for students.",
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.2,
    )

    content = response.choices[0].message.content

    # Try to parse the JSON directly
    try:
        cards = json.loads(content)
    except json.JSONDecodeError:
        # Try to recover by extracting the JSON array substring
        start = content.find("[")
        end = content.rfind("]")
        if start != -1 and end != -1 and start < end:
            cards = json.loads(content[start : end + 1])
        else:
            raise

    cleaned_cards = []
    for c in cards:
        if isinstance(c, dict):
            front = (c.get("front") or c.get("term") or "").strip()
            back = (c.get("back") or c.get("definition") or "").strip()
            if front and back:
                cleaned_cards.append({"front": front, "back": back})

    return cleaned_cards

def generate_quiz_from_text(
    text: str,
    num_questions: int = 10,
    difficulty: str = "mixed",
    include_mcq: bool = True,
    include_tf: bool = True,
) -> list[dict]:
    """
    Call the OpenAI API to generate quiz questions from the given text.

    Returns a list of dicts, each like:
    {
        "type": "mcq" or "true_false",
        "question": "...",
        "options": ["A", "B", "C", "D"] or ["True", "False"],
        "correct_answer": "A" or "True" etc,
        "difficulty": "easy"/"medium"/"hard",
        "explanation": "..."
    }
    """
    if not client:
        raise RuntimeError("OpenAI client is not configured. Missing OPENAI_API_KEY.")

    MAX_CHARS = 12000
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]

    # Describe which types to include for the prompt
    if include_mcq and include_tf:
        type_desc = "a mix of multiple-choice and true/false questions"
    elif include_mcq:
        type_desc = "multiple-choice questions only"
    elif include_tf:
        type_desc = "true/false questions only"
    else:
        # Should be prevented by the UI, but just in case:
        raise ValueError("At least one of include_mcq or include_tf must be True.")

    difficulty = difficulty.lower()
    if difficulty not in ["easy", "medium", "hard", "mixed"]:
        difficulty = "mixed"

    difficulty_instruction = (
        "Use a mix of easy, medium, and hard questions."
        if difficulty == "mixed"
        else f"Make all questions roughly {difficulty} difficulty."
    )

    prompt = f"""
You are an assistant that creates high-quality quiz questions to help a student study.

Given the course material below, generate up to {num_questions} {type_desc}.
{difficulty_instruction}

Output format (JSON):

[
  {{
    "type": "mcq" or "true_false",
    "question": "The question text",
    "options": ["option1", "option2", "option3", "option4"],  // For mcq only
    "correct_answer": "one of the options, or 'True'/'False' for true_false",
    "difficulty": "easy" or "medium" or "hard",
    "explanation": "1-3 sentence explanation of why the answer is correct"
  }},
  ...
]

Rules:
- For multiple-choice (mcq):
  - Provide exactly 4 options.
  - The correct_answer must match one of the options exactly.
- For true_false:
  - The options field can be omitted or set to ["True", "False"].
  - The correct_answer must be exactly "True" or "False".
- Focus on important concepts, formulas, and distinctions.
- Avoid trivial or overly obscure details.
- Respond with ONLY valid JSON, no extra commentary.

COURSE MATERIAL START
{text}
COURSE MATERIAL END
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You generate clear, accurate quiz questions for students.",
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.3,
    )

    content = response.choices[0].message.content

    try:
        questions = json.loads(content)
    except json.JSONDecodeError:
        # Try to recover by extracting the JSON array substring
        start = content.find("[")
        end = content.rfind("]")
        if start != -1 and end != -1 and start < end:
            questions = json.loads(content[start : end + 1])
        else:
            raise

    cleaned = []
    for q in questions:
        if not isinstance(q, dict):
            continue

        q_type = (q.get("type") or "").strip().lower()
        if q_type in ["mcq", "multiple_choice", "multiple-choice"]:
            q_type = "mcq"
        elif q_type in ["true_false", "true-false", "tf"]:
            q_type = "true_false"
        else:
            # Skip unknown types
            continue

        question_text = (q.get("question") or "").strip()
        if not question_text:
            continue

        options = q.get("options")
        if q_type == "mcq":
            # Ensure we have 4 options
            if not isinstance(options, list) or len(options) < 2:
                continue
            # Trim and keep up to 4
            options = [str(o).strip() for o in options if str(o).strip()]
            if len(options) < 2:
                continue
            options = options[:4]
        else:
            # true/false
            options = ["True", "False"]

        correct = (q.get("correct_answer") or "").strip()
        if q_type == "mcq":
            # Correct must match one of the options
            if correct not in options:
                # Try case-insensitive match
                lowered = [o.lower() for o in options]
                if correct.lower() in lowered:
                    correct = options[lowered.index(correct.lower())]
                else:
                    continue
        else:
            if correct not in ["True", "False"]:
                continue

        difficulty_val = (q.get("difficulty") or "").strip().lower()
        if difficulty_val not in ["easy", "medium", "hard"]:
            difficulty_val = "medium"

        explanation = (q.get("explanation") or "").strip()

        cleaned.append(
            {
                "type": q_type,
                "question": question_text,
                "options": options,
                "correct_answer": correct,
                "difficulty": difficulty_val,
                "explanation": explanation,
            }
        )

    return cleaned


def generate_cheatsheet_from_text(
    text: str,
    sheet_size: str = "3x5",
    focus: str = "both",
) -> str:
    """
    Call the OpenAI API to generate a compact cheat sheet from the given text.

    Returns a single markdown-formatted string that fits the requested size and focus.
    """
    if not client:
        raise RuntimeError("OpenAI client is not configured. Missing OPENAI_API_KEY.")

    # Trim text so the prompt doesn't get too huge
    MAX_CHARS = 16000
    if len(text) > MAX_CHARS:
        text = text[:MAX_CHARS]

    # Map sheet size to approximate max characters of output.
    # These are rough guidelines to force brevity.
    size = sheet_size.lower()
    if size == "3x5":
        max_chars_out = 800     # very compact
        size_desc = "a very compact 3x5 inch notecard"
    elif size == "1_page":
        max_chars_out = 2000
        size_desc = "a single-sided 8.5x11 inch page"
    else:  # "2_page"
        max_chars_out = 3500
        size_desc = "a two-sided 8.5x11 inch page"

    focus = focus.lower()
    if focus == "formulas":
        focus_desc = "ONLY key formulas and their very short labels. Avoid prose definitions."
    elif focus == "definitions":
        focus_desc = "key concepts and definitions, but not detailed derivations or long explanations."
    else:
        focus_desc = "a mix of key formulas, short definitions, and core concepts."

    prompt = f"""
You are an assistant that creates dense, exam-ready cheat sheets for students.

Given the course material below, create a cheat sheet that would fit on {size_desc}.
Focus on {focus_desc}

Constraints:
- Use bullet points and short, clear lines.
- Prioritize the most important ideas, formulas, and relationships.
- Avoid full sentences when a shorthand phrase works.
- Group related items with short headings when helpful.
- DO NOT include extra commentary about what you're doing.
- Your entire output must be at most {max_chars_out} characters.
- Output should be in plain text or simple markdown (headings + bullet points).

COURSE MATERIAL START
{text}
COURSE MATERIAL END
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You create concise, well-organized cheat sheets for exams.",
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.2,
    )

    content = response.choices[0].message.content

    # Extra safety: hard enforce the character limit
    if len(content) > max_chars_out:
        content = content[:max_chars_out]

    return content.strip()

def answer_question_with_materials(material_text: str, question: str) -> str:
    """
    Use the OpenAI API to answer a question based on the given course material text.
    The model is instructed to rely primarily on the provided material and to say
    when something is not covered.
    """
    if not client:
        raise RuntimeError("OpenAI client is not configured. Missing OPENAI_API_KEY.")

    # Keep context within a reasonable limit
    MAX_CHARS = 16000
    if len(material_text) > MAX_CHARS:
        material_text = material_text[:MAX_CHARS]

    prompt = f"""
You are a helpful study assistant for a university student.

You are given course materials (slides, notes, etc.) and then a question.
Your job is to answer the question as clearly and concisely as possible,
based primarily on the given materials.

If the materials do NOT provide enough information to fully answer,
say so explicitly and explain what is missing, rather than guessing.

COURSE MATERIALS:
-----------------
{material_text}

QUESTION:
---------
{question}
"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {
                "role": "system",
                "content": "You are a careful, concise tutor. Base your answers on the provided course materials.",
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.2,
    )

    answer = response.choices[0].message.content
    return answer.strip()



def strip_basic_markdown(text: str) -> str:
    """
    Remove simple markdown formatting markers like **bold**, __bold__, and inline `code`
    while keeping the inner text.
    """
    # bold: **text**
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    # bold via __text__
    text = re.sub(r"__(.+?)__", r"\1", text)
    # inline code: `text`
    text = text.replace("`", "")
    return text



def build_cheatsheet_pdf(title: str, body: str, sheet_size: str) -> bytes:
    """
    Build a PDF cheat sheet from the given title and body text.

    sheet_size: "3x5", "1_page", or "2_page"
    Returns PDF bytes.

    - 3x5: single small landscape notecard, 1 page max.
    - 1_page: letter page, 2 columns, 1 page max (extra text truncated).
    - 2_page: letter, 2 columns, up to 2 pages (extra text truncated).
    """
    buffer = io.BytesIO()

    # Choose page size
    if sheet_size == "3x5":
        # 5x3 inches, landscape notecard
        pagesize = (5 * inch, 3 * inch)
        max_pages = 1
        num_columns = 1
    else:
        # Standard 8.5x11 for both 1_page and 2_page
        pagesize = letter
        num_columns = 2
        max_pages = 1 if sheet_size == "1_page" else 2

    c = canvas.Canvas(buffer, pagesize=pagesize)
    width, height = pagesize

    # Layout parameters
    margin = 0.35 * inch if sheet_size == "3x5" else 0.75 * inch
    title_font_size = 10 if sheet_size == "3x5" else 14
    body_font_size = 7 if sheet_size == "3x5" else 10
    line_height = body_font_size + 2

    # Column settings
    if num_columns == 1:
        column_gap = 0.0
        column_width = width - 2 * margin
    else:
        column_gap = 0.3 * inch
        column_width = (width - 2 * margin - column_gap) / 2.0

    # We approximate how many characters fit in a line based on the column width.
    # This is rough but helps avoid overly long lines.
    avg_char_width = body_font_size * 0.55  # heuristic
    max_chars_per_line = max(25, int(column_width / avg_char_width))

    # Optional border so it visually looks like a card/page
    def draw_border():
        c.setLineWidth(1)
        c.rect(margin / 2, margin / 2, width - margin, height - margin)

    def start_page(page_number: int):
        c.setFont("Helvetica", body_font_size)
        draw_border()
        # Title on first page only
        if page_number == 1 and title.strip():
            c.setFont("Helvetica-Bold", title_font_size)
            c.drawString(margin, height - margin - title_font_size, title.strip())
            y_start = height - margin - title_font_size * 2
        else:
            y_start = height - margin
        c.setFont("Helvetica", body_font_size)
        return y_start

    # Prepare lines: parse simple markdown-style headings & bullets, wrap at word boundaries
    processed_lines = []
    for raw_line in body.splitlines():
        # Clean up markdown markers like **bold** so the PDF doesn't show the asterisks
        raw_line = strip_basic_markdown(raw_line)
        stripped = raw_line.strip()
        if not stripped:
            processed_lines.append("")
            continue

        # Headings: #, ##, ### -> uppercase, add a blank line before
        if stripped.startswith("### "):
            heading_text = stripped[4:].strip()
            if heading_text:
                processed_lines.append("")
                for wrap_line in textwrap.wrap(heading_text.upper(), width=max_chars_per_line):
                    processed_lines.append(wrap_line)
            continue
        elif stripped.startswith("## "):
            heading_text = stripped[3:].strip()
            if heading_text:
                processed_lines.append("")
                for wrap_line in textwrap.wrap(heading_text.upper(), width=max_chars_per_line):
                    processed_lines.append(wrap_line)
            continue
        elif stripped.startswith("# "):
            heading_text = stripped[2:].strip()
            if heading_text:
                processed_lines.append("")
                for wrap_line in textwrap.wrap(heading_text.upper(), width=max_chars_per_line):
                    processed_lines.append(wrap_line)
            continue

        # Bullets: "- " or "• " -> "• " with hanging indent
        if stripped.startswith("- ") or stripped.startswith("• "):
            bullet_text = stripped[2:].strip()
            wrapped = textwrap.wrap(bullet_text, width=max_chars_per_line) or [""]
            for i, w in enumerate(wrapped):
                if i == 0:
                    processed_lines.append("• " + w)
                else:
                    processed_lines.append("  " + w)
        else:
            # Normal text line, word-wrapped
            wrapped = textwrap.wrap(stripped, width=max_chars_per_line) or [""]
            processed_lines.extend(wrapped)


    # Drawing text in columns, respecting max_pages
    page_number = 1
    current_column = 0
    y = start_page(page_number)

    # Precompute column x positions
    column_x_positions = [
        margin + i * (column_width + column_gap) for i in range(num_columns)
    ]

    text_obj = c.beginText(column_x_positions[current_column], y)

    for line in processed_lines:
        # If out of vertical space, move to next column or next page
        if y <= margin:
            # Finish current column
            c.drawText(text_obj)

            current_column += 1
            if current_column < num_columns:
                # Next column on same page
                y = height - margin
                text_obj = c.beginText(column_x_positions[current_column], y)
            else:
                # Need a new page
                if page_number >= max_pages:
                    # We've reached the max allowed pages; stop drawing further lines
                    break
                c.showPage()
                page_number += 1
                draw_border()
                # Title only on first page; other pages start directly with body
                c.setFont("Helvetica", body_font_size)
                y = height - margin
                current_column = 0
                text_obj = c.beginText(column_x_positions[current_column], y)

        # Draw the line
        text_obj.textLine(line)
        y -= line_height

    # Draw remaining text object and finish
    c.drawText(text_obj)
    c.save()

    pdf_bytes = buffer.getvalue()
    buffer.close()
    return pdf_bytes


# -----------------------------------------
# Streamlit UI setup
# -----------------------------------------

ensure_data_dirs()

st.set_page_config(
    page_title="StudyPilot",
    page_icon="🎓",
    layout="wide",
)

st.title("StudyPilot 🎓")
st.write(
    "An AI-powered study companion that turns your course materials into a study plan, flashcards, quizzes, cheat sheets, and a Q&A chat."
)

# -----------------------------------------
# Sidebar state for course name input
# -----------------------------------------

if "course_name_input_version" not in st.session_state:
    st.session_state["course_name_input_version"] = 0

if "reset_course_name" not in st.session_state:
    st.session_state["reset_course_name"] = False

# If we requested a reset last run, bump the version so the widget key changes
if st.session_state["reset_course_name"]:
    st.session_state["course_name_input_version"] += 1
    st.session_state["reset_course_name"] = False

course_name_input_key = f"new_course_name_{st.session_state['course_name_input_version']}"

# -------------------------
# Sidebar: Course management
# -------------------------
st.sidebar.header("Course Selection & Management")

courses = load_courses()
course_names = [c["name"] for c in courses]

# Section to create a new course
st.sidebar.subheader("Create a new course")

new_course_name = st.sidebar.text_input(
    "New course name",
    key=course_name_input_key,
)

if st.sidebar.button("Add Course"):
    name = new_course_name.strip()
    if name == "":
        st.sidebar.error("Please enter a course name before adding.")
    else:
        existing = get_course_by_name(name)
        if existing:
            st.sidebar.warning("A course with that name already exists.")
        else:
            created = create_course(name)
            st.sidebar.success(f"Course '{created['name']}' created.")
            # Request the input to be cleared on next rerun
            st.session_state["reset_course_name"] = True
            st.rerun()

# Course selection
st.sidebar.subheader("Select a course")
if courses:
    selected_course_name = st.sidebar.selectbox(
        "Current course:",
        ["-- Select a course --"] + course_names,
        index=0,
    )
else:
    selected_course_name = st.sidebar.selectbox(
        "Current course:",
        ["No courses available"],
        index=0,
    )

selected_course = None
if courses and selected_course_name != "-- Select a course --":
    selected_course = get_course_by_name(selected_course_name)

if selected_course:
    st.sidebar.info(f"Selected course: **{selected_course['name']}**")

    # Delete course button
    if st.sidebar.button("Delete selected course"):
        delete_course(selected_course["id"])
        st.sidebar.success(f"Deleted course: {selected_course['name']}")
        st.rerun()
else:
    st.sidebar.info("No course selected yet. Create and select a course to enable uploads and features.")

# -------------------------
# Main tabs
# -------------------------
tabs = st.tabs(
    [
        "📂 Course Materials",
        "🗓️ Study Planner",
        "🃏 Flashcards",
        "📝 Quizzes",
        "📄 Cheat Sheets",
        "💬 Q&A Chat",
    ]
)

# 1) Course Materials tab
with tabs[0]:
    st.header("Course Materials")

    if not selected_course:
        st.warning("Please create and select a course in the sidebar to manage materials.")
    else:
        st.subheader(f"Files for: {selected_course['name']}")

        course_id = selected_course["id"]
        course_dir = get_course_dir(course_id)
        uploads_dir = course_dir / "uploads"
        uploads_dir.mkdir(parents=True, exist_ok=True)

        # -------------------------
        # Uploader state per course
        # -------------------------
        uploader_version_key = f"uploader_version_{course_id}"
        reset_uploader_flag_key = f"reset_uploader_{course_id}"

        if uploader_version_key not in st.session_state:
            st.session_state[uploader_version_key] = 0
        if reset_uploader_flag_key not in st.session_state:
            st.session_state[reset_uploader_flag_key] = False

        # If requested, bump version so uploader gets a new key and clears its files
        if st.session_state[reset_uploader_flag_key]:
            st.session_state[uploader_version_key] += 1
            st.session_state[reset_uploader_flag_key] = False

        uploader_key = f"uploader_{course_id}_{st.session_state[uploader_version_key]}"

        st.write("Upload lecture slides, notes, PDFs, or other course files here.")

        uploaded_files = st.file_uploader(
            "Choose files to upload",
            accept_multiple_files=True,
            key=uploader_key,
        )

        # Automatically process new uploads once, then clear uploader via version bump
        if uploaded_files:
            meta = load_course_meta(course_id)
            existing_files = meta.get("files", [])

            for uploaded in uploaded_files:
                stored_name = make_unique_filename(uploads_dir, uploaded.name)
                file_path = uploads_dir / stored_name

                # Save file to disk
                with open(file_path, "wb") as f:
                    f.write(uploaded.getbuffer())

                file_record = {
                    "original_name": uploaded.name,
                    "stored_name": stored_name,
                    "uploaded_at": datetime.now().isoformat(timespec="seconds"),
                    "size_bytes": len(uploaded.getbuffer()),
                }
                existing_files.append(file_record)

            meta["files"] = existing_files
            save_course_meta(course_id, meta)

            st.success(f"Uploaded {len(uploaded_files)} file(s) successfully.")

            # Clear the uploader selection on next run
            st.session_state[reset_uploader_flag_key] = True
            st.rerun()

        # List existing files
        meta = load_course_meta(course_id)
        files = meta.get("files", [])

        if not files:
            st.info("No files uploaded yet for this course.")
        else:
            st.write("### Uploaded Files")

            for i, f_info in enumerate(files):
                cols = st.columns([4, 2, 2, 1])
                cols[0].write(f"**{f_info['original_name']}**")
                cols[1].write(f_info["uploaded_at"])
                size_kb = f_info["size_bytes"] / 1024
                cols[2].write(f"{size_kb:.1f} KB")

                delete_button = cols[3].button(
                    "🗑️",
                    key=f"delete_{course_id}_{i}",
                    help="Delete this file",
                )

                if delete_button:
                    # Delete file from disk
                    file_path = uploads_dir / f_info["stored_name"]
                    if file_path.exists():
                        os.remove(file_path)

                    # Remove from meta
                    new_files = [ff for j, ff in enumerate(files) if j != i]
                    meta["files"] = new_files
                    save_course_meta(course_id, meta)
                    st.success(f"Deleted file: {f_info['original_name']}")

                    # Also clear the uploader selection so you don't have to click X
                    st.session_state[reset_uploader_flag_key] = True
                    st.rerun()

        st.info(
            "These uploaded files will be used later by the Study Planner, Flashcard Generator, Quiz Generator, "
            "Cheat Sheet Generator, and Q&A Chat."
        )

# 2) Study Planner tab
with tabs[1]:
    st.header("AI Study Schedule Planner")
    st.write(
        "Generate a high-level study schedule for a specific exam using the files you've uploaded for this course."
    )

    if not selected_course:
        st.warning("Please create and select a course in the sidebar first.")
    else:
        st.subheader(f"Study planner for: {selected_course['name']}")

        # Load files for this course
        course_id = selected_course["id"]
        meta = load_course_meta(course_id)
        files = meta.get("files", [])

        if not files:
            st.info(
                "No files uploaded yet for this course. "
                "Upload lecture slides or notes in the Course Materials tab first."
            )
        else:
            # Inputs for exam and dates
            col1, col2 = st.columns(2)
            with col1:
                exam_name = st.text_input("Exam name", value="Midterm 1")
                start_date = st.date_input(
                    "Start studying on",
                    value=date.today(),
                    help="This is the first day you plan to start studying for this exam.",
                )
            with col2:
                exam_date = st.date_input(
                    "Exam date",
                    value=date.today(),
                    help="You will study up to the day before this date.",
                )
                hours_per_day = st.slider(
                    "Target study hours per day for this course",
                    min_value=1.0,
                    max_value=8.0,
                    value=2.0,
                    step=0.5,
                )

            if exam_date <= start_date:
                st.error("The exam date must be after the study start date.")
            else:
                num_days = (exam_date - start_date).days  # we study up to exam_date - 1

                st.write(
                    f"Study window: **{start_date}** to **{exam_date}** "
                    f"(you have **{num_days}** day(s) to study before the exam)."
                )

                if st.button("Generate study plan"):
                    total_weight = sum(f["size_bytes"] for f in files)
                    if total_weight <= 0:
                        st.error("Unable to compute workload from files. Please re-upload or add files.")
                    else:
                        # Assign continuous date ranges to each file based on size
                        plan_rows = []
                        day_index = 0
                        num_files = len(files)

                        for idx, f_info in enumerate(files):
                            weight = f_info["size_bytes"]
                            share = weight / total_weight

                            # Approximate number of days for this file
                            days_remaining = num_days - day_index
                            if days_remaining <= 0:
                                file_days = 0
                            elif idx == num_files - 1:
                                # Last file gets all remaining days
                                file_days = days_remaining
                            else:
                                file_days = max(1, min(days_remaining, round(share * num_days)))

                            if file_days > 0:
                                start_d = start_date + timedelta(days=day_index)
                                end_d = start_date + timedelta(days=day_index + file_days - 1)

                                est_total_hours = share * num_days * hours_per_day

                                plan_rows.append(
                                    {
                                        "File": f_info["original_name"],
                                        "Suggested start date": str(start_d),
                                        "Suggested end date": str(end_d),
                                        "Estimated total hours": round(est_total_hours, 1),
                                    }
                                )

                                day_index += file_days

                        # If rounding left some unused days, extend the last file to the final day
                        if plan_rows and day_index < num_days:
                            last = plan_rows[-1]
                            last["Suggested end date"] = str(start_date + timedelta(days=num_days - 1))

                        st.subheader(f"File-level plan for: {exam_name}")
                        st.write(
                            "This table shows how the planner allocates your uploaded files across the available study days "
                            "based on their relative size."
                        )
                        if plan_rows:
                            st.table(plan_rows)
                        else:
                            st.info("No days available to schedule. Check your dates.")

                        # Build a simple day-by-day view
                        st.subheader("Day-by-day study schedule")
                        daily_rows = []
                        for offset in range(num_days):
                            day = start_date + timedelta(days=offset)
                            files_for_day = [
                                row["File"]
                                for row in plan_rows
                                if date.fromisoformat(row["Suggested start date"])
                                <= day
                                <= date.fromisoformat(row["Suggested end date"])
                            ]
                            daily_rows.append(
                                {
                                    "Date": str(day),
                                    "Planned study focus": ", ".join(files_for_day) if files_for_day else "—",
                                    "Suggested hours": hours_per_day if files_for_day else 0.0,
                                }
                            )

                        if daily_rows:
                            st.table(daily_rows)
                        else:
                            st.info("No schedule could be generated. Check your dates and files.")


# 3) Flashcards tab
with tabs[2]:
    st.header("AI Flashcard Generator")
    st.write(
        "Generate term/definition flashcards from your uploaded lecture slides and notes."
    )

    if not selected_course:
        st.warning("Please create and select a course in the sidebar first.")
    else:
        if client is None:
            st.error(
                "OpenAI API key not found. Please set OPENAI_API_KEY in a .env file "
                "before using the flashcard generator."
            )
        else:
            course_id = selected_course["id"]
            meta = load_course_meta(course_id)
            files = meta.get("files", [])

            if not files:
                st.info(
                    "No files uploaded yet for this course. "
                    "Upload lecture slides or notes in the Course Materials tab first."
                )
            else:
                st.subheader(f"Generate flashcards for: {selected_course['name']}")

                # Choose which files to include
                uploads_dir = get_course_dir(course_id) / "uploads"
                name_to_file = {f["original_name"]: f for f in files}

                default_selection = list(name_to_file.keys())
                selected_file_names = st.multiselect(
                    "Select files to include in this flashcard set",
                    options=list(name_to_file.keys()),
                    default=default_selection,
                )

                if not selected_file_names:
                    st.warning("Please select at least one file to generate flashcards from.")
                else:
                    num_cards = st.slider(
                        "Number of flashcards",
                        min_value=5,
                        max_value=100,
                        value=20,
                        step=5,
                    )

                    set_name = st.text_input(
                        "Flashcard set name",
                        value="Default set",
                        help="This name is just for display in this session.",
                    )

                    if st.button("Generate flashcards"):
                        # Aggregate text from selected files
                        all_text_parts = []
                        for fname in selected_file_names:
                            f_info = name_to_file[fname]
                            file_path = uploads_dir / f_info["stored_name"]
                            extracted = extract_text_from_file(file_path)
                            if extracted:
                                all_text_parts.append(extracted)

                        combined_text = "\n\n".join(all_text_parts).strip()

                        if not combined_text:
                            st.error(
                                "Could not extract text from the selected files. "
                                "Try different files or formats (PDF, PPTX, TXT)."
                            )
                        else:
                            with st.spinner("Generating flashcards with AI..."):
                                try:
                                    cards = generate_flashcards_from_text(combined_text, num_cards=num_cards)
                                except Exception as e:
                                    st.error(f"Error generating flashcards: {e}")
                                else:
                                    if not cards:
                                        st.warning(
                                            "The AI did not return any valid flashcards. "
                                            "Try reducing the number of files or simplifying the content."
                                        )
                                    else:
                                        # Store in session so they persist while app is running
                                        state_key = f"flashcards_{course_id}"
                                        st.session_state[state_key] = {
                                            "set_name": set_name,
                                            "cards": cards,
                                        }
                                        st.success(
                                            f"Generated {len(cards)} flashcard(s) for set: {set_name}"
                                        )

            # Show flashcards if we have a generated set in this session
            state_key = f"flashcards_{selected_course['id']}" if selected_course else None
            if state_key and state_key in st.session_state:
                flash_state = st.session_state[state_key]
                st.subheader(f"Flashcard set: {flash_state['set_name']}")
                st.write(
                    "Click each card to view the back. You can regenerate a new set at any time."
                )

                for idx, card in enumerate(flash_state["cards"], start=1):
                    with st.expander(f"Card {idx}: {card['front']}"):
                        st.write(card["back"])


# 4) Quizzes tab
with tabs[3]:
    st.header("AI Quiz Generator")
    st.write(
        "Generate multiple-choice and true/false questions from your uploaded course materials."
    )

    if not selected_course:
        st.warning("Please create and select a course in the sidebar first.")
    else:
        if client is None:
            st.error(
                "OpenAI API key not found. Please set OPENAI_API_KEY in a .env file "
                "before using the quiz generator."
            )
        else:
            course_id = selected_course["id"]
            meta = load_course_meta(course_id)
            files = meta.get("files", [])

            if not files:
                st.info(
                    "No files uploaded yet for this course. "
                    "Upload lecture slides or notes in the Course Materials tab first."
                )
            else:
                uploads_dir = get_course_dir(course_id) / "uploads"
                name_to_file = {f["original_name"]: f for f in files}

                st.subheader(f"Generate quiz for: {selected_course['name']}")

                selected_file_names = st.multiselect(
                    "Select files to include in this quiz",
                    options=list(name_to_file.keys()),
                    default=list(name_to_file.keys()),
                )

                if not selected_file_names:
                    st.warning("Please select at least one file to generate questions from.")
                else:
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        num_questions = st.slider(
                            "Approximate Number of Questions",
                            min_value=5,
                            max_value=40,
                            value=10,
                            step=1,
                        )
                    with col2:
                        difficulty_choice = st.selectbox(
                            "Difficulty",
                            options=["Mixed", "Easy", "Medium", "Hard"],
                            index=0,
                            help="You can choose a single difficulty or let the AI mix them.",
                        )
                    with col3:
                        include_mcq = st.checkbox("Include MCQ", value=True)
                        include_tf = st.checkbox("Include True/False", value=True)

                    if not include_mcq and not include_tf:
                        st.error("Please select at least one question type (MCQ or True/False).")
                    else:
                        if st.button("Generate quiz"):
                            # Aggregate text from selected files
                            all_text_parts = []
                            for fname in selected_file_names:
                                f_info = name_to_file[fname]
                                file_path = uploads_dir / f_info["stored_name"]
                                extracted = extract_text_from_file(file_path)
                                if extracted:
                                    all_text_parts.append(extracted)

                            combined_text = "\n\n".join(all_text_parts).strip()

                            if not combined_text:
                                st.error(
                                    "Could not extract text from the selected files. "
                                    "Try different files or formats (PDF, PPTX, TXT)."
                                )
                            else:
                                with st.spinner("Generating quiz with AI..."):
                                    try:
                                        questions = generate_quiz_from_text(
                                            combined_text,
                                            num_questions=num_questions,
                                            difficulty=difficulty_choice,
                                            include_mcq=include_mcq,
                                            include_tf=include_tf,
                                        )
                                    except Exception as e:
                                        st.error(f"Error generating quiz: {e}")
                                    else:
                                        if not questions:
                                            st.warning(
                                                "The AI did not return any valid questions. "
                                                "Try adjusting the number of questions or selected files."
                                            )
                                        else:
                                            state_key = f"quiz_{course_id}"
                                            st.session_state[state_key] = {
                                                "questions": questions,
                                                "meta": {
                                                    "num_questions": len(questions),
                                                    "difficulty": difficulty_choice,
                                                },
                                            }
                                            st.success(f"Generated {len(questions)} question(s).")

            # Display quiz questions if we have a set in session
            quiz_state_key = f"quiz_{selected_course['id']}" if selected_course else None
            if quiz_state_key and quiz_state_key in st.session_state:
                quiz_state = st.session_state[quiz_state_key]
                questions = quiz_state["questions"]

                st.subheader("Generated Quiz")
                st.write(
                    "Select an answer for each question. Feedback and explanations will appear after you choose."
                )

                course_id = selected_course["id"]
                total_answered = 0
                total_correct = 0

                for idx, q in enumerate(questions, start=1):
                    display_type = "MCQ"
                    if q["type"] == "true_false":
                        display_type = "T/F"

                    q_label = f"Q{idx} [{display_type} | {q['difficulty'].capitalize()}]"
                    st.markdown(f"### {q_label}")


                    # Options for this question
                    options = q["options"] if q["type"] == "mcq" else ["True", "False"]

                    # Unique key per question so selection is remembered
                    q_key = f"quiz_{course_id}_q_{idx}"

                    selected = st.radio(
                        "Your answer:",
                        options,
                        index=None,  # No default selection
                        key=q_key,
                    )

                    # Only evaluate once the user has picked something
                    if selected is not None:
                        total_answered += 1
                        if selected == q["correct_answer"]:
                            total_correct += 1
                            st.success("✅ Correct!")
                        else:
                            st.error(f"❌ Incorrect. Correct answer: **{q['correct_answer']}**")

                        if q["explanation"]:
                            st.info(f"Explanation: {q['explanation']}")

                    st.markdown("---")

                # Simple score summary
                st.subheader("Quiz Progress")
                st.write(
                    f"Answered **{total_answered}** out of **{len(questions)}** questions. "
                    f"Correct so far: **{total_correct}**."
                )



# 5) Cheat Sheets tab
with tabs[4]:
    st.header("AI Cheat Sheet Generator")
    st.write(
        "Build a compact cheat sheet (3×5 card, 1 page, or 2-sided) from your uploaded course materials."
    )

    if not selected_course:
        st.warning("Please create and select a course in the sidebar first.")
    else:
        if client is None:
            st.error(
                "OpenAI API key not found. Please set OPENAI_API_KEY in a .env file "
                "before using the cheat sheet generator."
            )
        else:
            course_id = selected_course["id"]
            meta = load_course_meta(course_id)
            files = meta.get("files", [])

            if not files:
                st.info(
                    "No files uploaded yet for this course. "
                    "Upload lecture slides or notes in the Course Materials tab first."
                )
            else:
                uploads_dir = get_course_dir(course_id) / "uploads"
                name_to_file = {f["original_name"]: f for f in files}

                st.subheader(f"Generate cheat sheet for: {selected_course['name']}")

                selected_file_names = st.multiselect(
                    "Select files to include",
                    options=list(name_to_file.keys()),
                    default=list(name_to_file.keys()),
                )

                if not selected_file_names:
                    st.warning("Please select at least one file to generate a cheat sheet from.")
                else:
                    col1, col2 = st.columns(2)
                    with col1:
                        sheet_size_label = st.selectbox(
                            "Cheat sheet size",
                            options=["3×5 Notecard", "1 Page", "2-Sided Page"],
                            index=1,
                            help="This affects how compact the cheat sheet will be.",
                        )
                    with col2:
                        focus_label = st.selectbox(
                            "Content focus",
                            options=["Formulas only", "Definitions/concepts", "Both"],
                            index=2,
                            help="Choose what type of content to emphasize.",
                        )

                    # Map labels to internal codes
                    size_map = {
                        "3×5 Notecard": "3x5",
                        "1 Page": "1_page",
                        "2-Sided Page": "2_page",
                    }
                    focus_map = {
                        "Formulas only": "formulas",
                        "Definitions/concepts": "definitions",
                        "Both": "both",
                    }

                    sheet_size = size_map[sheet_size_label]
                    focus = focus_map[focus_label]

                    custom_title = st.text_input(
                        "Optional cheat sheet title",
                        value="Exam Cheat Sheet",
                        help="This will appear at the top of the generated cheat sheet.",
                    )

                    if st.button("Generate cheat sheet"):
                        # Aggregate text from selected files
                        all_text_parts = []
                        for fname in selected_file_names:
                            f_info = name_to_file[fname]
                            file_path = uploads_dir / f_info["stored_name"]
                            extracted = extract_text_from_file(file_path)
                            if extracted:
                                all_text_parts.append(extracted)

                        combined_text = "\n\n".join(all_text_parts).strip()

                        if not combined_text:
                            st.error(
                                "Could not extract text from the selected files. "
                                "Try different files or formats (PDF, PPTX, TXT)."
                            )
                        else:
                            with st.spinner("Generating cheat sheet with AI..."):
                                try:
                                    body = generate_cheatsheet_from_text(
                                        combined_text,
                                        sheet_size=sheet_size,
                                        focus=focus,
                                    )
                                except Exception as e:
                                    st.error(f"Error generating cheat sheet: {e}")
                                else:
                                    # Build display text (Markdown-style)
                                    full_text = body
                                    if custom_title.strip():
                                        full_text = f"# {custom_title.strip()}\n\n" + body

                                    # Build PDF bytes matching the selected size
                                    pdf_bytes = build_cheatsheet_pdf(
                                        title=custom_title.strip() or "Exam Cheat Sheet",
                                        body=body,
                                        sheet_size=sheet_size,
                                    )

                                    state_key = f"cheatsheet_{course_id}"
                                    st.session_state[state_key] = {
                                        "title": custom_title.strip() or "Exam Cheat Sheet",
                                        "text": full_text,
                                        "pdf_bytes": pdf_bytes,
                                        "sheet_size": sheet_size_label,
                                    }
                                    st.success("Cheat sheet generated successfully.")


            # Display cheat sheet
            cheatsheet_state_key = f"cheatsheet_{selected_course['id']}" if selected_course else None
            if cheatsheet_state_key and cheatsheet_state_key in st.session_state:
                cs = st.session_state[cheatsheet_state_key]
                st.subheader("Generated Cheat Sheet")

                # Show markdown preview
                st.markdown(cs["text"])

                st.write(f"Format: **{cs['sheet_size']}**")

                # Download as PDF
                st.download_button(
                    label="Download cheat sheet as PDF",
                    data=cs["pdf_bytes"],
                    file_name=f"{cs['title'].replace(' ', '_')}.pdf",
                    mime="application/pdf",
                )

                # still offer plain text download
                st.download_button(
                    label="Download as .txt",
                    data=cs["text"],
                    file_name=f"{cs['title'].replace(' ', '_')}.txt",
                    mime="text/plain",
                )



# 6) Q&A Chat tab
with tabs[5]:
    st.header("Course Q&A Chat")
    st.write(
        "Ask questions about this course and get answers grounded in your uploaded materials."
    )

    if not selected_course:
        st.warning("Please create and select a course in the sidebar first.")
    else:
        if client is None:
            st.error(
                "OpenAI API key not found. Please set OPENAI_API_KEY in a .env file "
                "before using the Q&A chat."
            )
        else:
            course_id = selected_course["id"]
            meta = load_course_meta(course_id)
            files = meta.get("files", [])

            if not files:
                st.info(
                    "No files uploaded yet for this course. "
                    "Upload lecture slides or notes in the Course Materials tab first."
                )
            else:
                uploads_dir = get_course_dir(course_id) / "uploads"
                name_to_file = {f["original_name"]: f for f in files}

                st.subheader(f"Chat about: {selected_course['name']}")

                selected_file_names = st.multiselect(
                    "Select files to use as context for answers",
                    options=list(name_to_file.keys()),
                    default=list(name_to_file.keys()),
                    help="The assistant will base its answers primarily on these files.",
                )

                if not selected_file_names:
                    st.warning("Please select at least one file so the assistant has context.")
                else:
                    # Initialize chat history for this course
                    history_key = f"qa_history_{course_id}"
                    if history_key not in st.session_state:
                        st.session_state[history_key] = []

                    # Show existing chat history
                    st.write("### Chat")
                    for msg in st.session_state[history_key]:
                        if msg["role"] == "user":
                            with st.chat_message("user"):
                                st.markdown(msg["content"])
                        else:
                            with st.chat_message("assistant"):
                                st.markdown(msg["content"])

                    # Chat input
                    user_question = st.chat_input(
                        "Ask a question about this course (definitions, concepts, explanations, etc.)"
                    )

                    if user_question:
                        # Add user message to history
                        st.session_state[history_key].append(
                            {"role": "user", "content": user_question}
                        )

                        # Build material context from selected files
                        all_text_parts = []
                        for fname in selected_file_names:
                            f_info = name_to_file[fname]
                            file_path = uploads_dir / f_info["stored_name"]
                            extracted = extract_text_from_file(file_path)
                            if extracted:
                                all_text_parts.append(extracted)

                        combined_text = "\n\n".join(all_text_parts).strip()

                        if not combined_text:
                            assistant_reply = (
                                "I couldn't extract any readable text from the selected files. "
                                "Try different files or formats (PDF, PPTX, TXT)."
                            )
                        else:
                            try:
                                with st.spinner("Thinking..."):
                                    assistant_reply = answer_question_with_materials(
                                        combined_text, user_question
                                    )
                            except Exception as e:
                                assistant_reply = f"Error while answering: {e}"

                        # Add assistant response to history
                        st.session_state[history_key].append(
                            {"role": "assistant", "content": assistant_reply}
                        )

                        # Rerun so the new messages appear immediately
                        st.rerun()
